"""PptToHtmlConverter — converts PPTX/PPT files to HTML for ingestion.

Each slide becomes a <section class="slide"> element.
Text frames (title, content, bullet lists) are emitted as headings/paragraphs.
Tables are preserved as <table>/<tr>/<th>/<td>.
Images / diagrams are sent to the vision model (if configured) for a detailed
HTML description, with OCR as fallback.

Caching
-------
Converted HTML is saved to {output_dir}/{stem}.html and reused on subsequent
runs as long as the HTML file is newer than the source file.
"""

from __future__ import annotations

import html as _html_lib
import logging
from pathlib import Path
from typing import TYPE_CHECKING

from policygpt.ingestion.converters.base import HtmlConverter

if TYPE_CHECKING:
    from policygpt.ingestion.converters.vision import VisionDescriber
    from policygpt.ingestion.extraction.ocr import OcrExtractor

logger = logging.getLogger(__name__)


class PptToHtmlConverter(HtmlConverter):
    """Converts PPTX/PPT files to HTML.

    Parameters
    ----------
    output_dir:
        Directory where converted HTML files are written.
    skip_if_cached:
        Reuse existing HTML when newer than the source file.
    vision_describer:
        Optional LLM vision model for image/diagram shapes.
    ocr:
        Optional OCR extractor fallback for image shapes.
    """

    def __init__(
        self,
        output_dir: str | Path,
        skip_if_cached: bool = True,
        vision_describer: "VisionDescriber | None" = None,
        ocr: "OcrExtractor | None" = None,
    ) -> None:
        super().__init__(output_dir=output_dir, skip_if_cached=skip_if_cached)
        self._vision = vision_describer
        self._ocr = ocr

    @property
    def supported_content_types(self) -> frozenset[str]:
        return frozenset({"pptx", "ppt"})

    def _convert_to_html(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ImportError(
                "PPTX-to-HTML conversion requires 'python-pptx'. "
                "Install it with: pip install python-pptx"
            ) from exc

        title = self._title_from_stem(path.stem)
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        slide_parts: list[str] = []

        logger.info("PptToHtmlConverter: %s — %d slide(s)", path.name, slide_count)
        print(f"    [PPT] {path.name} — {slide_count} slide(s)", flush=True)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_html = self._convert_slide(slide, slide_num, path.name)
            if slide_html.strip():
                slide_parts.append(slide_html)

        meta = (
            f'<dl class="doc-meta">'
            f"<dt>Source</dt><dd>{_html_lib.escape(path.name)}</dd>"
            f"<dt>Type</dt><dd>PowerPoint</dd>"
            f"<dt>Slides</dt><dd>{slide_count}</dd>"
            f"</dl>"
        )
        body = meta + "\n" + "\n".join(slide_parts)
        return self._wrap_html(title, body)

    def _convert_slide(self, slide, slide_num: int, filename: str) -> str:
        """Convert one slide to an HTML <section> fragment."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        parts: list[str] = []
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            if shape.has_table:
                parts.append(self._table_to_html(shape.table))
            elif shape.has_text_frame:
                text_html = self._text_frame_to_html(shape)
                if text_html:
                    parts.append(text_html)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_html = self._describe_image_shape(shape, slide_num, filename)
                if img_html:
                    parts.append(img_html)

        if not parts:
            return ""
        inner = "\n".join(parts)
        return f'<section class="slide" data-slide="{slide_num}">\n{inner}\n</section>'

    def _describe_image_shape(self, shape, slide_num: int, filename: str) -> str:
        """Send a picture shape through vision → OCR → skip chain."""
        try:
            image_bytes: bytes = shape.image.blob
            mime_type: str = shape.image.content_type or "image/png"
        except Exception as exc:
            logger.warning(
                "PptToHtmlConverter: slide %d — failed to read image blob: %s", slide_num, exc
            )
            return ""

        size_kb = len(image_bytes) / 1024
        logger.info(
            "PptToHtmlConverter: slide %d image — %.1f KB (%s)", slide_num, size_kb, mime_type
        )
        print(
            f"    [PPT] slide {slide_num} image — {size_kb:.1f} KB ({mime_type})", flush=True
        )

        # 1. Vision model
        if self._vision is not None:
            logger.info(
                "PptToHtmlConverter: slide %d — vision [%s] …",
                slide_num, self._vision.provider_name,
            )
            print(
                f"    [PPT] slide {slide_num} — vision [{self._vision.provider_name}] …",
                flush=True,
            )
            html_fragment = self._vision.describe_page(image_bytes, mime_type)
            if html_fragment.strip():
                logger.info(
                    "PptToHtmlConverter: slide %d — vision OK (%d chars)", slide_num, len(html_fragment)
                )
                print(
                    f"    [PPT] slide {slide_num} — vision OK ({len(html_fragment)} chars)",
                    flush=True,
                )
                return (
                    f'<div class="slide-image" '
                    f'data-source="vision:{self._vision.provider_name}">\n'
                    f"{html_fragment}\n"
                    f"</div>"
                )
            logger.warning(
                "PptToHtmlConverter: slide %d — vision empty, trying OCR", slide_num
            )
            print(f"    [PPT] slide {slide_num} — vision empty, trying OCR …", flush=True)

        # 2. OCR fallback
        if self._ocr is not None:
            logger.info("PptToHtmlConverter: slide %d — OCR …", slide_num)
            print(f"    [PPT] slide {slide_num} — OCR …", flush=True)
            ocr_text = self._ocr.extract_from_bytes(image_bytes, mime_type)
            if ocr_text.strip():
                lines = [
                    f"<p>{_html_lib.escape(line)}</p>"
                    for line in ocr_text.splitlines()
                    if line.strip()
                ]
                logger.info(
                    "PptToHtmlConverter: slide %d — OCR OK (%d chars)", slide_num, len(ocr_text)
                )
                print(
                    f"    [PPT] slide {slide_num} — OCR OK ({len(ocr_text)} chars)", flush=True
                )
                return (
                    f'<div class="slide-image" data-source="ocr">\n'
                    f'{chr(10).join(lines)}\n'
                    f"</div>"
                )

        logger.warning(
            "PptToHtmlConverter: slide %d — no vision/OCR output, image skipped", slide_num
        )
        print(f"    [PPT] slide {slide_num} — image skipped (no vision/OCR output)", flush=True)
        return ""

    @staticmethod
    def _text_frame_to_html(shape) -> str:
        frame = shape.text_frame
        parts: list[str] = []

        for para_idx, para in enumerate(frame.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            escaped = _html_lib.escape(text)

            is_title_ph = shape.shape_type in (13, 14, 15)
            if para_idx == 0 and is_title_ph:
                parts.append(f"<h2>{escaped}</h2>")
            elif para.runs and any(
                r.font.bold for r in para.runs if r.font and r.font.bold is not None
            ):
                parts.append(f"<h3>{escaped}</h3>")
            elif para.level and para.level > 0:
                indent = "  " * para.level
                parts.append(f"{indent}<li>{escaped}</li>")
            else:
                parts.append(f"<p>{escaped}</p>")

        return "\n".join(parts)

    @staticmethod
    def _table_to_html(table) -> str:
        rows = list(table.rows)
        if not rows:
            return ""

        header_cells = "".join(
            f'<th scope="col">{_html_lib.escape(cell.text.strip()) or "&#8212;"}</th>'
            for cell in rows[0].cells
        )
        thead = f"<thead><tr>{header_cells}</tr></thead>"

        body_rows: list[str] = []
        for row in rows[1:]:
            cells = "".join(
                f"<td>{_html_lib.escape(cell.text.strip()) or '&#8212;'}</td>"
                for cell in row.cells
            )
            body_rows.append(f"<tr>{cells}</tr>")
        tbody = ("<tbody>\n" + "\n".join(body_rows) + "\n</tbody>") if body_rows else ""

        return f"<table>\n{thead}\n{tbody}\n</table>"
