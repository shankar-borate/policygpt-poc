"""
PolicyGPT Hybrid Search – Engineering Presentation Generator
Creates a professional PPTX with architecture diagrams, data-flow slides,
and deep technical details for the engineering team.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette (dark-tech theme) ─────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # near-black navy
CARD_BG      = RGBColor(0x16, 0x23, 0x3E)   # panel dark-blue
ACCENT_TEAL  = RGBColor(0x00, 0xD4, 0xC2)   # primary teal
ACCENT_BLUE  = RGBColor(0x38, 0x8B, 0xFF)   # secondary blue
ACCENT_AMBER = RGBColor(0xFF, 0xB8, 0x2C)   # highlight amber
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # success green
ACCENT_PURP  = RGBColor(0xA7, 0x8B, 0xFA)   # vector / semantic purple
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xBB, 0xC4, 0xD5)
MID_GREY     = RGBColor(0x55, 0x65, 0x80)
DIVIDER      = RGBColor(0x1E, 0x30, 0x55)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

blank_layout = prs.slide_layouts[6]   # completely blank


# ═══════════════════════════════════════════════════════════════════════════════
# Helper utilities
# ═══════════════════════════════════════════════════════════════════════════════

def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=Pt(1.5), corner=0.1):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, x, y, w, h, font_size=Pt(12), bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True,
                 italic=False, font_name="Segoe UI"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb

def add_multiline_text(slide, lines, x, y, w, h, font_size=Pt(11),
                       color=WHITE, bold_first=False, line_space=1.15,
                       font_name="Segoe UI", align=PP_ALIGN.LEFT):
    """lines = list of (text, bold, color_override)"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color

        if first:
            para = tf.paragraphs[0]
            first = False
        else:
            para = tf.add_paragraph()
        para.alignment = align
        para.space_after  = Pt(2)
        para.space_before = Pt(2)
        run = para.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = col
        run.font.name = font_name
    return txb

def slide_background(slide, color=DARK_BG):
    bg = add_rect(slide, 0, 0, 13.33, 7.50, fill_color=color)

def add_slide_label(slide, label, color=ACCENT_TEAL):
    """Small ALL-CAPS section label top-left."""
    add_text_box(slide, label, 0.35, 0.18, 5, 0.28,
                 font_size=Pt(9), bold=True, color=color,
                 font_name="Segoe UI Semibold")

def add_slide_title(slide, title, subtitle=None, y=0.45):
    add_text_box(slide, title, 0.35, y, 12.5, 0.6,
                 font_size=Pt(26), bold=True, color=WHITE,
                 font_name="Segoe UI Semibold")
    if subtitle:
        add_text_box(slide, subtitle, 0.35, y + 0.62, 12.5, 0.35,
                     font_size=Pt(13), color=LIGHT_GREY, font_name="Segoe UI")

def h_divider(slide, y, color=DIVIDER, x=0.35, w=12.6):
    add_rect(slide, x, y, w, 0.018, fill_color=color)

def add_badge(slide, text, x, y, w=1.4, h=0.28, fill=ACCENT_TEAL, text_color=DARK_BG):
    add_rounded_rect(slide, x, y, w, h, fill_color=fill, corner=0.25)
    add_text_box(slide, text, x, y+0.02, w, h-0.02,
                 font_size=Pt(9), bold=True, color=text_color,
                 align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

def add_arrow_right(slide, x, y, w=0.5, color=ACCENT_TEAL):
    """Simple right-pointing arrow using a line + triangle approximation."""
    # Horizontal line
    line = slide.shapes.add_connector(1,
        Inches(x), Inches(y + 0.13),
        Inches(x + w - 0.12), Inches(y + 0.13))
    line.line.color.rgb = color
    line.line.width = Pt(2)

def add_arrow_down(slide, x, y, h=0.35, color=ACCENT_TEAL):
    line = slide.shapes.add_connector(1,
        Inches(x), Inches(y),
        Inches(x), Inches(y + h))
    line.line.color.rgb = color
    line.line.width = Pt(2)

def footer(slide, text="PolicyGPT – Hybrid Search Architecture  |  Engineering Review"):
    add_rect(slide, 0, 7.25, 13.33, 0.25, fill_color=RGBColor(0x0A, 0x10, 0x1E))
    add_text_box(slide, text, 0.35, 7.26, 10, 0.22,
                 font_size=Pt(8), color=MID_GREY, font_name="Segoe UI")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Cover
# ═══════════════════════════════════════════════════════════════════════════════
def slide_cover():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)

    # Accent stripe left
    add_rect(s, 0, 0, 0.06, 7.50, fill_color=ACCENT_TEAL)

    # Gradient-like decorative bars
    for i, (col, alpha) in enumerate([(ACCENT_TEAL, 0.18), (ACCENT_BLUE, 0.12), (ACCENT_PURP, 0.08)]):
        add_rect(s, 13.33 - 2.5 + i*0.8, 0, 0.6, 7.50, fill_color=col)

    # Tag
    add_badge(s, "ENGINEERING DEEP DIVE", 0.5, 1.20, w=2.4, h=0.30,
              fill=ACCENT_TEAL, text_color=DARK_BG)

    # Main title
    add_text_box(s, "PolicyGPT Hybrid Search", 0.5, 1.70, 10, 1.0,
                 font_size=Pt(44), bold=True, color=WHITE,
                 font_name="Segoe UI Semibold")
    add_text_box(s, "Architecture, Data Flows & Domain Optimisation", 0.5, 2.72, 10, 0.55,
                 font_size=Pt(20), color=ACCENT_TEAL, font_name="Segoe UI Light")

    h_divider(s, 3.45, color=ACCENT_TEAL, x=0.5, w=9)

    add_text_box(s, "A complete engineering walkthrough of how ingestion, "
                    "vector indexing, hybrid retrieval\nand LLM answering combine "
                    "to deliver accurate, grounded answers over enterprise policy documents.",
                 0.5, 3.60, 10, 0.9, font_size=Pt(13), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Agenda chips
    chips = ["Hybrid Search Concepts", "System Architecture", "Ingestion Flow",
             "Retrieval Flow", "Score Blending", "Domain Optimisation"]
    for i, chip in enumerate(chips):
        cx = 0.5 + (i % 3) * 3.9
        cy = 4.80 + (i // 3) * 0.50
        add_rounded_rect(s, cx, cy, 3.6, 0.36, fill_color=CARD_BG,
                         line_color=ACCENT_TEAL, line_width=Pt(0.75))
        add_text_box(s, f"  {i+1}.  {chip}", cx+0.05, cy+0.03, 3.5, 0.32,
                     font_size=Pt(11), color=LIGHT_GREY, font_name="Segoe UI")

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Why Hybrid Search?  (Concept)
# ═══════════════════════════════════════════════════════════════════════════════
def slide_concept():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "01  CONCEPT")
    add_slide_title(s, "Why Hybrid Search?",
                    "No single retrieval method dominates all query types")
    h_divider(s, 1.25)

    # Three method cards
    cards = [
        ("🔑", "BM25 Keyword", ACCENT_BLUE,
         "Sparse term matching\n\n"
         "• Exact clause numbers (Cl. 4.2.1)\n"
         "• Defined policy terms & acronyms\n"
         "• Names, dates, rule identifiers\n\n"
         "Strength: Zero ambiguity on exact matches\n"
         "Weakness: No understanding of meaning"),
        ("🔄", "More-Like-This", ACCENT_AMBER,
         "Vocabulary overlap\n\n"
         "• Paraphrases & synonyms\n"
         "• Policy jargon vs plain language\n"
         "• Similar sentence structures\n\n"
         "Strength: Handles linguistic variation\n"
         "Weakness: Misses pure semantic intent"),
        ("🧠", "Vector / kNN", ACCENT_PURP,
         "Dense semantic embedding\n\n"
         "• Intent-based queries (what do I win?)\n"
         "• Conceptually similar paragraphs\n"
         "• Cross-lingual or vague phrasing\n\n"
         "Strength: Captures true meaning\n"
         "Weakness: Can miss exact terminology"),
    ]
    for i, (icon, title, col, body) in enumerate(cards):
        cx = 0.35 + i * 4.32
        add_rounded_rect(s, cx, 1.40, 4.1, 4.90, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.5))
        # Colour top accent bar
        add_rect(s, cx, 1.40, 4.1, 0.08, fill_color=col)
        add_text_box(s, icon, cx + 0.15, 1.52, 0.6, 0.5, font_size=Pt(24))
        add_text_box(s, title, cx + 0.65, 1.57, 3.3, 0.40,
                     font_size=Pt(15), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        lines = []
        for j, line in enumerate(body.split("\n")):
            bold = (j == 0)
            c = WHITE if not bold else LIGHT_GREY
            lines.append((line, bold, c))
        add_multiline_text(s, lines, cx + 0.18, 2.05, 3.82, 4.10,
                           font_size=Pt(10.5), color=LIGHT_GREY)

    # Bottom verdict banner
    add_rounded_rect(s, 0.35, 6.42, 12.63, 0.72, fill_color=ACCENT_TEAL,
                     line_color=None)
    add_text_box(s, "✦  Hybrid blends all three simultaneously — weak signals in one channel "
                    "are compensated by strong signals in another, delivering higher recall AND precision.",
                 0.55, 6.48, 12.3, 0.60,
                 font_size=Pt(12), bold=True, color=DARK_BG, font_name="Segoe UI Semibold")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – System Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════
def slide_architecture():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "02  ARCHITECTURE")
    add_slide_title(s, "System Architecture Overview",
                    "Three decoupled layers — Ingestion, Index, Retrieval — with a FastAPI serving layer")
    h_divider(s, 1.25)

    # Layer titles
    layers = [
        ("INGESTION PIPELINE",   ACCENT_AMBER, 0.35),
        ("OPENSEARCH INDICES",   ACCENT_TEAL,  4.60),
        ("RETRIEVAL & ANSWER",   ACCENT_BLUE,  8.85),
    ]
    for title, col, x in layers:
        add_rounded_rect(s, x, 1.32, 3.85, 0.30, fill_color=col)
        add_text_box(s, title, x+0.05, 1.34, 3.75, 0.26,
                     font_size=Pt(9), bold=True, color=DARK_BG,
                     align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")

    # ── Ingestion column ──────────────────────────────────────────────────────
    ing_boxes = [
        ("📄  Raw Documents", "HTML / PDF / TXT / OCR images", CARD_BG),
        ("✂️  Chunker", "min 300 / target 1800 / max 3200 chars", CARD_BG),
        ("🏷️  Metadata Extractor", "type · version · date · audiences · tags", CARD_BG),
        ("🧠  LLM Summariser", "Section & doc summaries via Bedrock/GPT", CARD_BG),
        ("📐  Embedder", "Amazon Titan embed-text-v2  (1536-dim)", CARD_BG),
        ("❓  FAQ Generator", "Up to 30 Q/A pairs per document", CARD_BG),
    ]
    for i, (title, sub, col) in enumerate(ing_boxes):
        by = 1.72 + i * 0.80
        add_rounded_rect(s, 0.38, by, 3.80, 0.68, fill_color=col,
                         line_color=ACCENT_AMBER, line_width=Pt(0.75))
        add_text_box(s, title, 0.52, by+0.05, 3.55, 0.28,
                     font_size=Pt(10.5), bold=True, color=WHITE,
                     font_name="Segoe UI Semibold")
        add_text_box(s, sub, 0.52, by+0.33, 3.55, 0.28,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
        if i < len(ing_boxes) - 1:
            add_arrow_down(s, 2.28, by + 0.68, h=0.12, color=ACCENT_AMBER)

    # ── Index column ──────────────────────────────────────────────────────────
    idx_boxes = [
        ("_sections index", "section_id · doc_id · raw_text · summary\n"
                            "embedding (knn_vector 1536d) · metadata_tags\n"
                            "audiences · user_ids · domain",
         CARD_BG, ACCENT_TEAL),
        ("_documents index", "doc_id · title · type · version · effective_date\n"
                             "keywords · audiences · summary · user_ids",
         CARD_BG, ACCENT_BLUE),
        ("_faqs index",      "faq_id · question · answer · question_embedding\n"
                             "user_ids · domain",
         CARD_BG, ACCENT_PURP),
        ("_threads index",   "thread_id · user_id · recent_messages\n"
                             "display_messages · conversation_summary",
         CARD_BG, ACCENT_AMBER),
    ]
    for i, (title, sub, col, lc) in enumerate(idx_boxes):
        by = 1.72 + i * 1.32
        add_rounded_rect(s, 4.63, by, 3.80, 1.18, fill_color=col,
                         line_color=lc, line_width=Pt(1.0))
        add_text_box(s, f"policygpt{title}", 4.77, by+0.06, 3.55, 0.30,
                     font_size=Pt(10.5), bold=True, color=lc,
                     font_name="Segoe UI Semibold")
        add_text_box(s, sub, 4.77, by+0.38, 3.55, 0.72,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")

    # ── Retrieval column ──────────────────────────────────────────────────────
    ret_boxes = [
        ("⚡  FAQ Fast-path", "kNN on question_embedding ≥ 0.92 → return immediately\n(zero full-RAG cost)", ACCENT_GREEN),
        ("🔀  HybridSearcher", "3 concurrent searches → min-max normalise → weighted blend\nKeyword 35% + MLT 15% + kNN 50%", ACCENT_TEAL),
        ("📊  Reranker", "Semantic · lexical · title · metadata weight scoring\nPrunes candidates → top sections per doc", ACCENT_BLUE),
        ("🤖  LLM Answer", "Evidence + recent_chat injected into prompt\nBedrock 120B / Claude / GPT-5", ACCENT_PURP),
        ("📝  Thread Store", "recent_messages (LLM context)\ndisplay_messages (UI history, accumulated in OS)", ACCENT_AMBER),
    ]
    for i, (title, sub, lc) in enumerate(ret_boxes):
        by = 1.72 + i * 0.98
        add_rounded_rect(s, 8.88, by, 4.10, 0.84, fill_color=CARD_BG,
                         line_color=lc, line_width=Pt(0.75))
        add_text_box(s, title, 9.02, by+0.06, 3.85, 0.30,
                     font_size=Pt(10.5), bold=True, color=lc,
                     font_name="Segoe UI Semibold")
        add_text_box(s, sub, 9.02, by+0.38, 3.85, 0.40,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
        if i < len(ret_boxes) - 1:
            add_arrow_down(s, 10.93, by + 0.84, h=0.14, color=lc)

    # Horizontal arrows between layers
    add_arrow_right(s, 4.18, 3.95, w=0.45, color=ACCENT_AMBER)
    add_arrow_right(s, 8.43, 3.95, w=0.45, color=ACCENT_TEAL)

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Ingestion Data Flow
# ═══════════════════════════════════════════════════════════════════════════════
def slide_ingestion_flow():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "03  INGESTION DATA FLOW")
    add_slide_title(s, "Ingestion Pipeline — Step by Step",
                    "Documents flow through extraction → enrichment → embedding → indexing, async in background")
    h_divider(s, 1.25)

    steps = [
        ("1", "File Discovery", ACCENT_AMBER,
         "Glob patterns: *.html *.htm *.pdf *.txt\n"
         "Dedup: skip if source_path already in OS\n"
         "Skip: filenames containing '_summary'"),
        ("2", "Text Extraction", ACCENT_AMBER,
         "HTML → BeautifulSoup tag-aware parse\n"
         "PDF  → pdfminer text + OCR fallback\n"
         "TXT  → direct read\n"
         "Images → AWS Textract (OCR, conf ≥ 80%)"),
        ("3", "Section Chunking", ACCENT_BLUE,
         "Target: 1800 chars  |  Min: 300  |  Max: 3200\n"
         "Preserves heading hierarchy\n"
         "Assigns order_index per section"),
        ("4", "Metadata Extraction", ACCENT_BLUE,
         "LLM extracts: document_type, version,\n"
         "effective_date, audiences, keywords,\n"
         "metadata_tags, entity map"),
        ("5", "Summarisation", ACCENT_TEAL,
         "Section-level: each section → concise summary\n"
         "Document-level: map-reduce over section summaries\n"
         "Recursive merging when input > token budget"),
        ("6", "Embedding", ACCENT_TEAL,
         "Model: amazon.titan-embed-text-v2 (1536-dim)\n"
         "Input: section summary (not raw text)\n"
         "L2-normalised for cosine similarity"),
        ("7", "FAQ Generation", ACCENT_PURP,
         "LLM generates up to 30 Q/A pairs per doc\n"
         "Each question separately embedded\n"
         "Stored in _faqs index for fast-path lookup"),
        ("8", "OS Indexing", ACCENT_PURP,
         "_sections: full section payload + embedding\n"
         "_documents: doc-level metadata + summary\n"
         "_faqs: Q/A + question embeddings\n"
         "user_ids & domain tag on every record"),
    ]

    cols = 4
    for i, (num, title, col, body) in enumerate(steps):
        row = i // cols
        c   = i % cols
        bx = 0.35 + c * 3.24
        by = 1.40 + row * 2.82

        add_rounded_rect(s, bx, by, 3.10, 2.60, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.2))
        # Step number badge
        add_rounded_rect(s, bx + 0.12, by + 0.10, 0.38, 0.38, fill_color=col)
        add_text_box(s, num, bx + 0.12, by + 0.10, 0.38, 0.38,
                     font_size=Pt(13), bold=True, color=DARK_BG,
                     align=PP_ALIGN.CENTER, font_name="Segoe UI Semibold")
        add_text_box(s, title, bx + 0.58, by + 0.14, 2.40, 0.32,
                     font_size=Pt(11.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, bx + 0.14, by + 0.58, 2.88, 1.90,
                     font_size=Pt(9.5), color=LIGHT_GREY, font_name="Segoe UI")

        # Arrow between steps in same row
        if c < cols - 1:
            add_arrow_right(s, bx + 3.10, by + 1.20, w=0.14, color=col)

    # Bottom async note
    add_rounded_rect(s, 0.35, 7.00, 12.63, 0.30, fill_color=CARD_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(0.5))
    add_text_box(s, "⚡  Background async: server status → 'ready' immediately after bot creation; "
                    "ingestion runs concurrently; retrieval is fully independent.",
                 0.50, 7.02, 12.3, 0.26, font_size=Pt(9.5), color=ACCENT_TEAL,
                 font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Retrieval Data Flow
# ═══════════════════════════════════════════════════════════════════════════════
def slide_retrieval_flow():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "04  RETRIEVAL DATA FLOW")
    add_slide_title(s, "Retrieval Pipeline — Query to Answer",
                    "Every user message goes through analysis → hybrid search → rerank → LLM answer")
    h_divider(s, 1.25)

    # Main flow boxes (top lane)
    flow = [
        ("User Query", WHITE, CARD_BG, ACCENT_TEAL,
         "POST /api/chat\nthread_id + message\nuser_id filter applied"),
        ("Query Analyser", "⚙️", CARD_BG, ACCENT_BLUE,
         "context_dependent?\nintents · topic_hints\nexpanded_terms\ncanonical_question"),
        ("FAQ Fast-Path", "⚡", CARD_BG, ACCENT_GREEN,
         "kNN on _faqs index\ncos-sim ≥ 0.92?\n→ return immediately\n(skips RAG entirely)"),
        ("Embed Query", "📐", CARD_BG, ACCENT_PURP,
         "Titan embed-text-v2\n1536-dim vector\nOriginal question only\n(not expanded terms)"),
        ("Hybrid Search", "🔀", CARD_BG, ACCENT_TEAL,
         "3 concurrent searches\non _sections index\nUser filter pre-applied\n→ blended candidates"),
        ("Reranker", "📊", CARD_BG, ACCENT_AMBER,
         "Weight multi-signal score\nSemantic+lexical+title\n+metadata+parent\nTop N sections"),
        ("LLM Prompt", "🤖", CARD_BG, ACCENT_BLUE,
         "Evidence + summaries\nrecent_chat context\nconversation_summary\n→ grounded answer"),
    ]

    bw = 1.70
    gap = 0.11
    total = len(flow) * bw + (len(flow) - 1) * gap
    start_x = (13.33 - total) / 2

    for i, (title, icon, bg, lc, detail) in enumerate(flow):
        bx = start_x + i * (bw + gap)
        by = 1.38

        add_rounded_rect(s, bx, by, bw, 2.10, fill_color=bg,
                         line_color=lc, line_width=Pt(1.5))
        # Top colour bar
        add_rect(s, bx, by, bw, 0.06, fill_color=lc)
        add_text_box(s, title, bx + 0.08, by + 0.10, bw - 0.16, 0.30,
                     font_size=Pt(10), bold=True, color=lc,
                     font_name="Segoe UI Semibold")
        add_text_box(s, detail, bx + 0.08, by + 0.45, bw - 0.16, 1.55,
                     font_size=Pt(8.5), color=LIGHT_GREY, font_name="Segoe UI")

        if i < len(flow) - 1:
            ax = bx + bw
            add_arrow_right(s, ax, by + 0.95, w=gap + 0.05, color=lc)

    # ── Detail panel: Hybrid search internals ─────────────────────────────────
    add_rounded_rect(s, 0.35, 3.70, 12.63, 2.85, fill_color=CARD_BG,
                     line_color=DIVIDER, line_width=Pt(0.5))
    add_text_box(s, "HYBRID SEARCH INTERNALS", 0.55, 3.76, 5, 0.25,
                 font_size=Pt(9), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")

    # Three search channels
    channels = [
        ("BM25 Keyword", ACCENT_BLUE,
         "multi_match: best_fields\n"
         "Fields: section_title^4 · document_title^3\n"
         "        keywords^3 · summary^2 · raw_text^1\n"
         "Fuzziness: AUTO  |  Operator: OR\n"
         "Policy analyzer: lowercase+stop+snowball\n"
         "→ Raw BM25 score (unbounded)"),
        ("More-Like-This", ACCENT_AMBER,
         "more_like_this on 3 fields:\n"
         "raw_text · summary · section_title\n"
         "min_term_freq=1  min_doc_freq=1\n"
         "max_query_terms=25\n"
         "minimum_should_match=30%\n"
         "→ TF-IDF vocabulary overlap score"),
        ("Vector kNN", ACCENT_PURP,
         "HNSW approximate nearest-neighbour\n"
         "Space: cosinesimil  |  Engine: lucene\n"
         "ef_construction=128  M=16\n"
         "1536-dim Titan embedding\n"
         "k = max(top_k × 2, rerank_candidates)\n"
         "→ Cosine similarity [0, 1]"),
    ]
    for i, (title, col, body) in enumerate(channels):
        cx = 0.50 + i * 4.20
        add_rounded_rect(s, cx, 4.10, 3.95, 2.25, fill_color=DARK_BG,
                         line_color=col, line_width=Pt(1.0))
        add_text_box(s, title, cx + 0.15, 4.15, 3.65, 0.28,
                     font_size=Pt(10.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, cx + 0.15, 4.48, 3.65, 1.75,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")

    # Blend formula box
    add_text_box(s, "Blend Formula:",
                 0.50, 6.60, 2.20, 0.28, font_size=Pt(9.5),
                 bold=True, color=ACCENT_TEAL, font_name="Segoe UI Semibold")
    add_text_box(s,
                 "score = 0.35 × norm(BM25)  +  0.15 × norm(MLT)  +  0.50 × norm(cosine)",
                 2.70, 6.60, 8.50, 0.28, font_size=Pt(10),
                 color=WHITE, font_name="Segoe UI Semibold")
    add_text_box(s, "Each type min-max normalised to [0,1] independently before blending.",
                 0.50, 6.90, 10, 0.25, font_size=Pt(9),
                 color=MID_GREY, font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Index Schema Deep-Dive
# ═══════════════════════════════════════════════════════════════════════════════
def slide_index_schema():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "05  INDEX SCHEMA")
    add_slide_title(s, "OpenSearch Index Schema",
                    "Four purpose-built indices — each optimised for a different access pattern")
    h_divider(s, 1.25)

    # Sections index (largest, most detail)
    add_rounded_rect(s, 0.35, 1.35, 6.20, 5.65, fill_color=CARD_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(1.5))
    add_rect(s, 0.35, 1.35, 6.20, 0.07, fill_color=ACCENT_TEAL)
    add_text_box(s, "policygpt_sections  (PRIMARY RETRIEVAL INDEX)", 0.52, 1.42, 5.90, 0.30,
                 font_size=Pt(11), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")

    sections_fields = [
        ("section_id / doc_id", "keyword", LIGHT_GREY),
        ("document_title / section_title", "text  [policy_analyzer]  + keyword sub-field", LIGHT_GREY),
        ("raw_text / summary", "text  [policy_analyzer]  — full searchable content", LIGHT_GREY),
        ("embedding", "knn_vector dim=1536  HNSW cosinesimil  ef=128 M=16", ACCENT_PURP),
        ("section_type / metadata_tags / keywords", "keyword  — facet filtering & boosting", LIGHT_GREY),
        ("audiences", "keyword  — audience-based access control", LIGHT_GREY),
        ("user_ids", "keyword[]  — permission pre-filter at query time", ACCENT_AMBER),
        ("domain", "keyword  — domain isolation (policy / contest)", LIGHT_GREY),
        ("order_index", "integer  — section order within document", LIGHT_GREY),
        ("source_path", "keyword  — file path for document viewer link", LIGHT_GREY),
    ]
    for i, (field, ftype, col) in enumerate(sections_fields):
        by = 1.82 + i * 0.46
        add_text_box(s, field, 0.50, by, 2.60, 0.36,
                     font_size=Pt(9.5), bold=True, color=WHITE,
                     font_name="Segoe UI Semibold")
        add_text_box(s, ftype, 3.15, by, 3.20, 0.36,
                     font_size=Pt(9), color=col, font_name="Segoe UI")

    # Analyser box
    add_rounded_rect(s, 0.35, 6.48, 6.20, 0.42, fill_color=DARK_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(0.5))
    add_text_box(s, "policy_analyzer:  standard tokenizer  →  lowercase  →  stop words  →  snowball stemmer",
                 0.52, 6.52, 5.90, 0.30, font_size=Pt(9), color=ACCENT_TEAL,
                 font_name="Segoe UI")

    # Right column: three smaller indices
    right_indices = [
        ("policygpt_documents", ACCENT_BLUE,
         "doc_id · title · document_type · version\n"
         "effective_date · keywords · metadata_tags\n"
         "audiences · summary · user_ids · domain\n\n"
         "Use: document-level lookup, cache check\n"
         "     document_indexed_for_path() skip\n"
         "     search_documents() Google-style UI search"),
        ("policygpt_faqs", ACCENT_PURP,
         "faq_id · doc_id · document_title\n"
         "question · answer · source_path\n"
         "question_embedding knn_vector 1536d\n"
         "user_ids · domain\n\n"
         "Use: FAQ fast-path (cos ≥ 0.92 bypass RAG)\n"
         "     aggregate query evidence (top-30 FAQ hits)"),
        ("policygpt_threads", ACCENT_AMBER,
         "thread_id · user_id · title\n"
         "recent_messages (JSON) → LLM context\n"
         "display_messages (JSON) → accumulated UI\n"
         "conversation_summary · active_doc_ids\n"
         "created_at · updated_at\n\n"
         "Use: conversation persistence across restarts"),
    ]
    for i, (name, col, body) in enumerate(right_indices):
        by = 1.35 + i * 2.00
        add_rounded_rect(s, 6.70, by, 6.25, 1.85, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_rect(s, 6.70, by, 6.25, 0.06, fill_color=col)
        add_text_box(s, name, 6.85, by + 0.08, 5.95, 0.28,
                     font_size=Pt(10.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, 6.85, by + 0.42, 5.95, 1.30,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Score Blending Deep-Dive
# ═══════════════════════════════════════════════════════════════════════════════
def slide_score_blending():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "06  SCORE BLENDING")
    add_slide_title(s, "Hybrid Score Blending — In Detail",
                    "Min-max normalisation → weighted sum → final ranking used by the reranker")
    h_divider(s, 1.25)

    # Step-by-step formula walkthrough
    formula_steps = [
        ("Step 1 — Run 3 searches concurrently", ACCENT_TEAL,
         "ThreadPoolExecutor(max_workers=3) dispatches keyword_search, similarity_search, "
         "vector_search in parallel.\nEach call applies the user_id pre-filter at the OS level "
         "before scoring — no unauthorised sections ever appear.\n"
         "Over-fetch multiplier: top_k = max(requested × 2, rerank_section_candidates) "
         "to give the reranker enough candidates."),
        ("Step 2 — Min-max normalise each result set independently", ACCENT_AMBER,
         "BM25 scores are unbounded (e.g. 0 – 22). Cosine similarity is in [0,1]. MLT varies.\n"
         "For each type:  norm_score(x) = (x – min) / (max – min)    [span = 1 if min == max]\n"
         "Sections appearing in only some result sets receive 0.0 for missing types — "
         "they are not penalised beyond that."),
        ("Step 3 — Weighted blend", ACCENT_PURP,
         "final_score = w_keyword × norm_keyword  +  w_similarity × norm_similarity  +  w_vector × norm_vector\n\n"
         "Default weights:  keyword=0.30  similarity=0.20  vector=0.50\n"
         "Policy domain:    keyword=0.35  similarity=0.15  vector=0.50   (exact terms matter more)\n"
         "Contest domain:   keyword=0.20  similarity=0.15  vector=0.65   (intent queries dominate)\n"
         "Weights are normalised at runtime — they do not need to sum to 1.0."),
        ("Step 4 — Downstream reranker re-scores top candidates", ACCENT_BLUE,
         "Within each document, sections are re-scored using a multi-signal formula:\n"
         "  section_score = 0.36×semantic + 0.24×lexical + 0.16×parent_doc + 0.12×title + 0.12×metadata\n"
         "  doc_score     = 0.48×semantic + 0.24×lexical + 0.16×title      + 0.12×metadata\n"
         "Final top_docs documents and top_sections_per_doc sections passed to LLM evidence context."),
    ]

    for i, (title, col, body) in enumerate(formula_steps):
        by = 1.40 + i * 1.44
        add_rounded_rect(s, 0.35, by, 12.63, 1.30, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.2))
        add_rect(s, 0.35, by, 0.06, 1.30, fill_color=col)
        add_text_box(s, title, 0.55, by + 0.06, 12.0, 0.28,
                     font_size=Pt(11), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, 0.55, by + 0.38, 12.0, 0.85,
                     font_size=Pt(9.5), color=LIGHT_GREY, font_name="Segoe UI")

    # Worked example
    add_rounded_rect(s, 0.35, 7.04, 12.63, 0.36, fill_color=DARK_BG,
                     line_color=ACCENT_GREEN, line_width=Pt(0.5))
    add_text_box(s,
                 "Example: Section A: BM25=18 (norm=1.0), MLT=0 (norm=0.0), cosine=0.71 (norm=0.85)  →  "
                 "0.35×1.0 + 0.15×0.0 + 0.50×0.85 = 0.775",
                 0.52, 7.07, 12.2, 0.28, font_size=Pt(9.5), color=ACCENT_GREEN,
                 font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Access Control & Permission Model
# ═══════════════════════════════════════════════════════════════════════════════
def slide_access_control():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "07  ACCESS CONTROL")
    add_slide_title(s, "Permission Model & Access Control",
                    "user_id pre-filters every search at the index level — no post-hoc filtering needed")
    h_divider(s, 1.25)

    # Flow: ingestion side
    add_rounded_rect(s, 0.35, 1.38, 5.90, 2.60, fill_color=CARD_BG,
                     line_color=ACCENT_AMBER, line_width=Pt(1.2))
    add_text_box(s, "INGESTION — Permission Assignment", 0.52, 1.44, 5.60, 0.28,
                 font_size=Pt(10.5), bold=True, color=ACCENT_AMBER,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "1.  Config.ingestion_user_ids  (env: POLICY_GPT_INGESTION_USER_IDS)\n"
                 "    Comma-separated list of user IDs granted access to ALL documents.\n\n"
                 "2.  Every section and document record stored with:\n"
                 "    user_ids: [\"100\", \"200\", \"301\"]  (keyword array)\n\n"
                 "3.  FAQs inherit user_ids from their parent document.\n\n"
                 "4.  Threads store user_id for listing and isolation.",
                 0.52, 1.80, 5.60, 2.10, font_size=Pt(10), color=LIGHT_GREY,
                 font_name="Segoe UI")

    add_arrow_right(s, 6.25, 2.55, w=0.55, color=ACCENT_AMBER)

    # Flow: retrieval side
    add_rounded_rect(s, 6.80, 1.38, 6.18, 2.60, fill_color=CARD_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(1.2))
    add_text_box(s, "RETRIEVAL — Permission Enforcement", 6.97, 1.44, 5.90, 0.28,
                 font_size=Pt(10.5), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "1.  user_id resolved from: URL ?user_id= param → cookie → empty.\n\n"
                 "2.  If hybrid_search_enabled=True and user_id is empty → HTTP 401.\n\n"
                 "3.  Every search query wraps the scoring clause in a bool/filter:\n"
                 "    { \"filter\": [ {\"terms\": {\"user_ids\": [\"100\"]}} ] }\n"
                 "    Applied BEFORE scoring — zero-cost, OS-native.\n\n"
                 "4.  Result: only sections the user is entitled to ever surface.",
                 6.97, 1.80, 5.90, 2.10, font_size=Pt(10), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Code-style DSL example
    add_rounded_rect(s, 0.35, 4.10, 12.63, 2.00, fill_color=RGBColor(0x0A, 0x10, 0x1E),
                     line_color=ACCENT_TEAL, line_width=Pt(0.8))
    add_text_box(s, "OpenSearch Query (simplified)", 0.52, 4.15, 5, 0.25,
                 font_size=Pt(9), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 '{\n'
                 '  "query": {\n'
                 '    "bool": {\n'
                 '      "must":   { "knn": { "embedding": { "vector": [...1536 floats...], "k": 24 } } },\n'
                 '      "filter": [ { "terms": { "user_ids": ["100"] } } ]\n'
                 '    }\n'
                 '  }\n'
                 '}',
                 0.52, 4.42, 12.10, 1.60, font_size=Pt(10),
                 color=ACCENT_GREEN, font_name="Consolas")

    # Notes
    notes = [
        ("Threads API", ACCENT_AMBER,
         "GET /api/threads?user_id=100 returns only threads owned by user 100.\n"
         "No cross-user thread leakage possible."),
        ("Search UI", ACCENT_BLUE,
         "POST /api/search enforces same user_id filter.\n"
         "BM25 multi_match also wrapped in bool/filter before scoring."),
        ("FAQ fast-path", ACCENT_PURP,
         "kNN on _faqs index also uses terms filter on user_ids.\n"
         "Even cache hits are permission-scoped."),
    ]
    for i, (title, col, body) in enumerate(notes):
        bx = 0.35 + i * 4.32
        add_rounded_rect(s, bx, 6.22, 4.10, 0.98, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(0.75))
        add_text_box(s, title, bx + 0.15, 6.27, 3.80, 0.25,
                     font_size=Pt(10), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, bx + 0.15, 6.56, 3.80, 0.58,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Domain Optimisation
# ═══════════════════════════════════════════════════════════════════════════════
def slide_domain_optimisation():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "08  DOMAIN OPTIMISATION")
    add_slide_title(s, "Domain Optimisation",
                    "Single config switch adapts the entire stack — weights, prompts, token budgets, OCR")
    h_divider(s, 1.25)

    # Domain comparison table
    add_rounded_rect(s, 0.35, 1.35, 12.63, 3.30, fill_color=CARD_BG,
                     line_color=DIVIDER, line_width=Pt(0.5))

    # Header row
    headers = ["Parameter", "Default", "policy  domain", "contest  domain", "Why it differs"]
    col_x = [0.40, 2.55, 4.45, 6.55, 8.65]
    col_w = [2.10, 1.85, 2.05, 2.05, 4.45]
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_text_box(s, h, x, 1.42, w, 0.28, font_size=Pt(10),
                     bold=True, color=ACCENT_TEAL,
                     font_name="Segoe UI Semibold")

    rows = [
        ("keyword_weight",    "0.30",  "0.35",  "0.20",
         "Policy: clause numbers & defined terms need exact match"),
        ("similarity_weight", "0.20",  "0.15",  "0.15",
         "Both domains: MLT adds marginal value over vector"),
        ("vector_weight",     "0.50",  "0.50",  "0.65",
         "Contest: vague user intent → semantic search dominates"),
        ("chat_max_tokens",   "2700",  "5400",  "2700",
         "Policy answers: multi-step procedures need room"),
        ("ocr_enabled",       "False", "True",  "True",
         "Both: scanned images, org charts, approval matrices"),
        ("ocr_min_conf",      "80%",   "80%",   "80%",
         "AWS Textract LINE block confidence threshold"),
        ("accuracy_profile",  "high",  "vhigh", "high",
         "Policy needs broader retrieval (vhigh: top 6 docs, 15 sections)"),
    ]
    row_colors = [LIGHT_GREY, LIGHT_GREY, ACCENT_PURP,
                  ACCENT_AMBER, ACCENT_GREEN, LIGHT_GREY, ACCENT_TEAL]
    for i, (row, col) in enumerate(zip(rows, row_colors)):
        ry = 1.80 + i * 0.38
        if i % 2 == 0:
            add_rect(s, 0.38, ry - 0.02, 12.55, 0.36, fill_color=DARK_BG)
        for j, (val, x, w) in enumerate(zip(row, col_x, col_w)):
            tc = col if j == 0 else (ACCENT_TEAL if j == 2 else (ACCENT_AMBER if j == 3 else LIGHT_GREY))
            add_text_box(s, val, x, ry, w, 0.34, font_size=Pt(9.5),
                         color=tc, font_name="Segoe UI")

    # Accuracy profiles
    add_text_box(s, "ACCURACY PROFILES", 0.40, 4.80, 3, 0.28,
                 font_size=Pt(10), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")

    profiles = [
        ("vhigh", ACCENT_TEAL,
         "top_docs=3  sections/doc=5  max_to_llm=15\n"
         "rerank_candidates=15  evidence_neighbors=5\n"
         "doc_summary_budget=8000 tokens\n"
         "→ Best quality, highest cost"),
        ("high", ACCENT_BLUE,
         "top_docs=3  sections/doc=3  max_to_llm=4\n"
         "rerank_candidates=12  evidence_neighbors=1\n"
         "doc_summary_budget=6000 tokens\n"
         "→ Balanced quality/cost"),
        ("medium", ACCENT_AMBER,
         "top_docs=2  sections/doc=2  max_to_llm=3\n"
         "rerank_candidates=8  evidence_neighbors=1\n"
         "→ Cost-optimised"),
        ("low", ACCENT_PURP,
         "top_docs=1  sections/doc=2  max_to_llm=2\n"
         "rerank_candidates=6\n"
         "→ Minimum viable retrieval"),
    ]
    for i, (name, col, body) in enumerate(profiles):
        bx = 0.35 + i * 3.20
        add_rounded_rect(s, bx, 5.12, 3.05, 2.10, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_badge(s, name, bx + 0.12, 5.20, w=1.10, h=0.28, fill=col, text_color=DARK_BG)
        add_text_box(s, body, bx + 0.14, 5.57, 2.80, 1.55,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Conversation Management
# ═══════════════════════════════════════════════════════════════════════════════
def slide_conversation():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "09  CONVERSATION MANAGEMENT")
    add_slide_title(s, "Conversation State & Thread Persistence",
                    "Two message stores per thread — lean in-memory LLM context + full OS-persisted display history")
    h_divider(s, 1.25)

    # Two-column layout
    # Left: message types
    add_rounded_rect(s, 0.35, 1.38, 5.90, 4.00, fill_color=CARD_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(1.2))
    add_text_box(s, "Two Message Stores per Thread", 0.52, 1.44, 5.60, 0.28,
                 font_size=Pt(11), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")

    add_rounded_rect(s, 0.52, 1.85, 5.55, 1.50, fill_color=DARK_BG,
                     line_color=ACCENT_BLUE, line_width=Pt(0.8))
    add_text_box(s, "recent_messages  (LLM context window)", 0.65, 1.90, 5.30, 0.25,
                 font_size=Pt(10), bold=True, color=ACCENT_BLUE,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "• Kept in memory at all times — never cleared\n"
                 "• Trimmed to last max_recent_messages=6 turns\n"
                 "• Injected into every LLM answer prompt as 'Recent:'\n"
                 "• Serialised to OS so survives server restart",
                 0.65, 2.18, 5.30, 1.10, font_size=Pt(9.5), color=LIGHT_GREY,
                 font_name="Segoe UI")

    add_rounded_rect(s, 0.52, 3.47, 5.55, 1.50, fill_color=DARK_BG,
                     line_color=ACCENT_PURP, line_width=Pt(0.8))
    add_text_box(s, "display_messages  (UI history)", 0.65, 3.52, 5.30, 0.25,
                 font_size=Pt(10), bold=True, color=ACCENT_PURP,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "• Accumulated in OS across all turns (not overwritten)\n"
                 "• Cleared from memory after each OS save (lean)\n"
                 "• Loaded fresh from OS on each API response\n"
                 "• Full message history visible to user on page refresh",
                 0.65, 3.80, 5.30, 1.10, font_size=Pt(9.5), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Right: turn-by-turn state machine
    add_rounded_rect(s, 6.45, 1.38, 6.53, 4.00, fill_color=CARD_BG,
                     line_color=ACCENT_AMBER, line_width=Pt(1.2))
    add_text_box(s, "Turn-by-Turn State Machine", 6.62, 1.44, 6.20, 0.28,
                 font_size=Pt(11), bold=True, color=ACCENT_AMBER,
                 font_name="Segoe UI Semibold")

    turns = [
        ("Turn N — chat()", ACCENT_AMBER,
         "get_thread() → in-memory (recent_msgs intact)\n"
         "Build prompt with recent_chat history\n"
         "Append [user: Q, asst: A] to both msg stores"),
        ("save_thread()", ACCENT_TEAL,
         "Load existing OS display_messages\n"
         "Prepend to new messages → accumulate\n"
         "Save full state to OS\n"
         "Clear display_messages from memory"),
        ("get_thread_for_display()", ACCENT_BLUE,
         "Load from OS → full display_messages\n"
         "Return to API → serialise to frontend\n"
         "User sees complete conversation history"),
    ]
    for i, (title, col, body) in enumerate(turns):
        by = 1.85 + i * 1.12
        add_rounded_rect(s, 6.62, by, 6.18, 0.98, fill_color=DARK_BG,
                         line_color=col, line_width=Pt(0.75))
        add_text_box(s, f"{i+1}. {title}", 6.78, by + 0.06, 5.88, 0.25,
                     font_size=Pt(10), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, 6.78, by + 0.34, 5.88, 0.58,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
        if i < len(turns) - 1:
            add_arrow_down(s, 9.71, by + 0.98, h=0.14, color=col)

    # Bottom: summarisation
    add_rounded_rect(s, 0.35, 5.50, 12.63, 1.30, fill_color=CARD_BG,
                     line_color=ACCENT_GREEN, line_width=Pt(1.0))
    add_text_box(s, "Auto Conversation Summarisation  (after summarize_after_turns=8 messages)",
                 0.52, 5.56, 12.0, 0.25, font_size=Pt(10.5), bold=True,
                 color=ACCENT_GREEN, font_name="Segoe UI Semibold")
    add_text_box(s,
                 "When len(recent_messages) ≥ 8  →  LLM generates a compact conversation_summary.\n"
                 "The summary is injected into the answer prompt as 'Summary:' (before 'Recent:') — "
                 "enabling coherent long conversations that would otherwise overflow the context window.\n"
                 "recent_messages is NOT cleared after summarisation — the two stores evolve independently.",
                 0.52, 5.85, 12.0, 0.88, font_size=Pt(9.5), color=LIGHT_GREY,
                 font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – FAQ Fast-Path
# ═══════════════════════════════════════════════════════════════════════════════
def slide_faq():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "10  FAQ FAST-PATH")
    add_slide_title(s, "FAQ Fast-Path — Zero-Cost Answer Cache",
                    "Pre-generated Q/A pairs bypass the full RAG pipeline when the query is a near-exact match")
    h_divider(s, 1.25)

    # Diagram
    # Ingestion side
    add_rounded_rect(s, 0.35, 1.40, 3.80, 3.30, fill_color=CARD_BG,
                     line_color=ACCENT_AMBER, line_width=Pt(1.2))
    add_text_box(s, "Ingestion — FAQ Generation", 0.52, 1.46, 3.60, 0.28,
                 font_size=Pt(10.5), bold=True, color=ACCENT_AMBER,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "1.  LLM generates up to 30 Q/A pairs\n"
                 "    per document during ingestion.\n\n"
                 "2.  Each question is embedded\n"
                 "    (Titan embed-text-v2, 1536-dim).\n\n"
                 "3.  Stored in policygpt_faqs index:\n"
                 "    question_embedding (knn_vector)\n"
                 "    question + answer text\n"
                 "    user_ids + domain filter\n\n"
                 "4.  Represents the 'canonical' Q/A\n"
                 "    knowledge base for the corpus.",
                 0.52, 1.80, 3.60, 2.80, font_size=Pt(10), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Arrow
    add_arrow_right(s, 4.15, 2.95, w=0.55, color=ACCENT_AMBER)

    # Fast-path decision
    add_rounded_rect(s, 4.70, 1.40, 4.20, 3.30, fill_color=CARD_BG,
                     line_color=ACCENT_GREEN, line_width=Pt(1.5))
    add_text_box(s, "⚡  Fast-Path Lookup at Query Time", 4.87, 1.46, 3.95, 0.28,
                 font_size=Pt(10.5), bold=True, color=ACCENT_GREEN,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "1.  Embed the raw user question.\n\n"
                 "2.  kNN search on question_embedding\n"
                 "    with user_id filter (k=1).\n\n"
                 "3.  If top score ≥ 0.92:\n"
                 "    ✓  Return stored FAQ answer\n"
                 "    ✓  Skip retrieval, reranker, LLM\n"
                 "    ✓  Fastest possible response\n\n"
                 "4.  If score < 0.92:\n"
                 "    →  Fall through to full RAG pipeline.",
                 4.87, 1.80, 3.95, 2.80, font_size=Pt(10), color=LIGHT_GREY,
                 font_name="Segoe UI")

    add_arrow_right(s, 8.90, 2.95, w=0.55, color=ACCENT_GREEN)

    # Aggregate FAQ use
    add_rounded_rect(s, 9.45, 1.40, 3.53, 3.30, fill_color=CARD_BG,
                     line_color=ACCENT_PURP, line_width=Pt(1.2))
    add_text_box(s, "Aggregate Query Evidence", 9.62, 1.46, 3.28, 0.28,
                 font_size=Pt(10.5), bold=True, color=ACCENT_PURP,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "For aggregate intents\n"
                 "(e.g. 'list all schemes'):\n\n"
                 "• Top-30 FAQ hits searched\n"
                 "  across ALL documents\n\n"
                 "• Grouped by document title\n"
                 "  for clean prompt layout\n\n"
                 "• Replaces section evidence\n"
                 "  for listing questions\n\n"
                 "• Gives corpus-wide coverage\n"
                 "  beyond doc-level retrieval.",
                 9.62, 1.80, 3.28, 2.80, font_size=Pt(10), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Cost savings callout
    add_rounded_rect(s, 0.35, 4.85, 12.63, 1.30, fill_color=CARD_BG,
                     line_color=ACCENT_GREEN, line_width=Pt(1.0))
    add_text_box(s, "Cost & Latency Impact", 0.52, 4.90, 5, 0.25,
                 font_size=Pt(10.5), bold=True, color=ACCENT_GREEN,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "A FAQ fast-path hit eliminates:  (1) retrieval query + embedding call  "
                 "(2) HybridSearcher + 3 OS queries  (3) Reranker computation  "
                 "(4) LLM chat call (largest cost).\n"
                 "For high-frequency, repetitive enterprise queries (eligibility, deadlines, limits) "
                 "the fast-path hit rate can exceed 40%, delivering substantial inference cost savings.",
                 0.52, 5.20, 12.0, 0.88, font_size=Pt(9.5), color=LIGHT_GREY,
                 font_name="Segoe UI")

    # Threshold note
    add_rounded_rect(s, 0.35, 6.28, 12.63, 0.72, fill_color=DARK_BG,
                     line_color=ACCENT_AMBER, line_width=Pt(0.5))
    add_text_box(s,
                 "Threshold  faq_fastpath_min_score=0.92  is tunable in config. "
                 "Lowering to 0.88 increases hit-rate but risks returning slightly mismatched answers. "
                 "Raising to 0.95 makes fast-path very conservative — only near-verbatim matches qualify.",
                 0.52, 6.32, 12.10, 0.60, font_size=Pt(9.5), color=LIGHT_GREY,
                 font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – LLM Prompt Construction
# ═══════════════════════════════════════════════════════════════════════════════
def slide_prompt():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "11  LLM PROMPT CONSTRUCTION")
    add_slide_title(s, "LLM Answer Prompt — Structure & Context Injection",
                    "Evidence-first architecture: every claim the LLM makes must be grounded in retrieved text")
    h_divider(s, 1.25)

    # Prompt structure diagram
    sections_prompt = [
        ("Summary:", ACCENT_GREEN,
         "conversation_summary (if exists)\nCompact rolling summary after 8 turns", 0.62),
        ("Recent:", ACCENT_TEAL,
         "Last 6 turns of recent_messages\n"
         "USER: ... / ASSISTANT: ...\nProvides turn-by-turn context continuity", 0.78),
        ("Question:", ACCENT_BLUE,
         "canonical_question from QueryAnalyser\n"
         "+ expanded_terms + topic_hints", 0.52),
        ("Style:", ACCENT_AMBER,
         "Answer format guidance per intent\n"
         "(exact / broad / aggregate / list)", 0.52),
        ("Docs:", ACCENT_PURP,
         "Aliases D1=, D2= → document titles\n"
         "Metadata: type, tags, audience, version", 0.52),
        ("Evidence:", ACCENT_TEAL,
         "Top sections: raw_text snippets\n"
         "+ section summary + section_type\n"
         "Neighboring units included (evidence_neighbors=5)", 0.78),
        ("Answer instruction:", ACCENT_GREEN,
         "\"Answer only from the evidence. Include exclusions/exceptions.\n"
         "Never mention internal IDs. Prefer raw evidence over summaries.\"", 0.65),
    ]

    bx = 0.35
    by = 1.38
    for i, (label, col, body, h) in enumerate(sections_prompt):
        add_rounded_rect(s, bx, by, 3.60, h, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_rect(s, bx, by, 0.06, h, fill_color=col)
        add_text_box(s, label, bx + 0.12, by + 0.03, 1.10, h - 0.06,
                     font_size=Pt(10), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, bx + 1.25, by + 0.03, 2.30, h - 0.06,
                     font_size=Pt(8.5), color=LIGHT_GREY, font_name="Segoe UI")
        by += h + 0.07

    # Right panel: query analysis detail
    add_rounded_rect(s, 4.15, 1.38, 8.83, 5.65, fill_color=CARD_BG,
                     line_color=ACCENT_BLUE, line_width=Pt(1.2))
    add_text_box(s, "QueryAnalyser — What Drives the Prompt", 4.32, 1.44, 8.50, 0.28,
                 font_size=Pt(11), bold=True, color=ACCENT_BLUE,
                 font_name="Segoe UI Semibold")

    qa_fields = [
        ("original_question",   "Raw user text — used for embedding (not expanded)", LIGHT_GREY),
        ("canonical_question",  "Normalised, cleaned question for BM25", LIGHT_GREY),
        ("context_dependent",   "True if follow-up (uses active_doc_ids, conversation topic)", ACCENT_AMBER),
        ("intents",             "exact_lookup | broad_lookup | aggregate | list | ...", LIGHT_GREY),
        ("topic_hints",         "Inferred policy domain tags → metadata pre-filter boost", ACCENT_GREEN),
        ("expanded_terms",      "Synonyms + acronym expansions for BM25 recall lift", LIGHT_GREY),
        ("exact_match_expected","True → evidece mode = exact (tighter char limits)", ACCENT_PURP),
        ("sub_questions",       "Compound question detection → separate retrieval per sub-Q", ACCENT_TEAL),
    ]
    for i, (field, desc, col) in enumerate(qa_fields):
        ry = 1.85 + i * 0.56
        if i % 2 == 0:
            add_rect(s, 4.18, ry - 0.04, 8.73, 0.52, fill_color=DARK_BG)
        add_text_box(s, field, 4.32, ry, 2.30, 0.40,
                     font_size=Pt(9.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, desc, 6.65, ry, 6.15, 0.40,
                     font_size=Pt(9.5), color=LIGHT_GREY, font_name="Segoe UI")

    # Intent → retrieval mode table
    add_rounded_rect(s, 4.15, 6.25, 8.83, 0.82, fill_color=DARK_BG,
                     line_color=ACCENT_TEAL, line_width=Pt(0.5))
    add_text_box(s, "Intent → Retrieval Mode",
                 4.30, 6.28, 2.50, 0.25, font_size=Pt(9), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")
    add_text_box(s,
                 "exact_lookup → top 2 docs × 4 sections × 3 to LLM  |  "
                 "broad_lookup → top 6 docs × 6 sections × 8 to LLM  |  "
                 "aggregate → FAQ evidence (top-30 hits cross-corpus)  |  "
                 "context_dependent → reuses active_doc_ids from prior turn",
                 4.30, 6.56, 8.55, 0.44, font_size=Pt(9), color=LIGHT_GREY,
                 font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Deployment & Configuration
# ═══════════════════════════════════════════════════════════════════════════════
def slide_deployment():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "12  DEPLOYMENT & CONFIGURATION")
    add_slide_title(s, "Deployment & Configuration Reference",
                    "All knobs exposed as environment variables — zero code change for domain, model, or scale switches")
    h_divider(s, 1.25)

    # Config groups
    config_groups = [
        ("AI / Model", ACCENT_BLUE, [
            ("ai_profile",        "openai | bedrock-20b | bedrock-120b | bedrock-claude-sonnet-4-6"),
            ("accuracy_profile",  "vhigh | high | medium | low"),
            ("runtime_cost_profile", "standard | aggressive"),
            ("bedrock_region",    "ap-south-1  (AWS Bedrock region)"),
        ]),
        ("Domain", ACCENT_AMBER, [
            ("domain_type",       "policy | contest  (selects prompt, weights, OCR)"),
            ("ingestion_user_ids","Comma-separated user IDs, e.g. 100,200,301"),
            ("document_folder",   "Path to policy document directory"),
            ("ocr_enabled",       "True/False  (overridden per domain)"),
        ]),
        ("OpenSearch", ACCENT_TEAL, [
            ("OS_HOST",           "OpenSearch cluster hostname"),
            ("OS_PORT",           "9200 (default)"),
            ("OS_USERNAME / OS_PASSWORD", "Credentials (load from env/secrets manager)"),
            ("OS_INDEX_PREFIX",   "policygpt  (renames all 4 indices)"),
            ("OS_USE_SSL / OS_VERIFY_CERTS", "True/False"),
        ]),
        ("Hybrid Search Weights", ACCENT_PURP, [
            ("hybrid_keyword_weight",    "0.30 default | 0.35 policy | 0.20 contest"),
            ("hybrid_similarity_weight", "0.20 default | 0.15 both domains"),
            ("hybrid_vector_weight",     "0.50 default | 0.50 policy | 0.65 contest"),
        ]),
        ("Retrieval Tuning", ACCENT_GREEN, [
            ("faq_fastpath_min_score", "0.92  (lower = more fast-path hits)"),
            ("max_recent_messages",    "6 (turns kept in LLM context window)"),
            ("summarize_after_turns",  "8 (auto-summarise long conversations)"),
            ("top_docs / top_sections_per_doc", "3 / 5 (vhigh) – see accuracy_profile"),
        ]),
    ]

    # Layout: 2×2 grid for first 4 groups, last group spans bottom row
    layout = [
        (0.35,  1.38, 6.25, 2.10),  # row0 col0
        (6.75,  1.38, 6.23, 2.10),  # row0 col1
        (0.35,  3.58, 6.25, 2.10),  # row1 col0
        (6.75,  3.58, 6.23, 2.10),  # row1 col1
        (0.35,  5.78, 12.63, 1.22), # row2 full-width
    ]
    for i, ((bx, by, bw, bh), (title, col, items)) in enumerate(zip(layout, config_groups)):
        add_rounded_rect(s, bx, by, bw, bh, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_text_box(s, title, bx + 0.15, by + 0.07, bw - 0.3, 0.26,
                     font_size=Pt(10.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        item_row_h = (bh - 0.44) / max(len(items), 1)
        for j, (key, val) in enumerate(items):
            ky = by + 0.40 + j * item_row_h
            # full-width last group: side-by-side inline
            if i == 4:
                val_x = bx + 0.18 + (bw / len(items)) * j
                add_text_box(s, f"{key}: {val}", val_x, ky, bw/len(items) - 0.1, item_row_h,
                             font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
            else:
                add_text_box(s, key, bx + 0.18, ky, 2.50, item_row_h,
                             font_size=Pt(9), bold=True, color=WHITE,
                             font_name="Segoe UI Semibold")
                add_text_box(s, val, bx + 2.72, ky, bw - 2.85, item_row_h,
                             font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")

    if False:  # dead code placeholder removed
        bx = 0
        bw = 0
        by = 0
        bh = 0
        add_rounded_rect(s, bx, by, bw, bh, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_text_box(s, title, bx + 0.15, by + 0.05, bw - 0.3, 0.24,
                     font_size=Pt(10), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        body = "  |  ".join(f"{k}: {v}" for k, v in items)
        add_text_box(s, body, bx + 0.18, by + 0.33, bw - 0.30, 0.52,
                     font_size=Pt(8.5), color=LIGHT_GREY, font_name="Segoe UI")

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 – Key Metrics & Observability
# ═══════════════════════════════════════════════════════════════════════════════
def slide_metrics():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_slide_label(s, "13  OBSERVABILITY")
    add_slide_title(s, "Observability & Key Metrics",
                    "Built-in usage tracking, retrieval logs, and health endpoints for monitoring")
    h_divider(s, 1.25)

    # Metric cards top row
    metric_cards = [
        ("LLM Usage Tracker", ACCENT_AMBER,
         "Tracks per-request:\n• Input tokens + cost (USD & INR)\n• Output tokens + cost\n"
         "• Rolling totals across session\n• Per-model pricing snapshot\n"
         "• Exposed at GET /api/usage\n• Live widget in UI sidebar"),
        ("Health Endpoint", ACCENT_TEAL,
         "GET /api/health returns:\n• status: ready|ingesting|error\n"
         "• document_count (live OS count)\n• section_count (live OS count)\n"
         "• thread_count (in-memory)\n• ingestion progress %\n• current_file being processed"),
        ("Retrieval Debug Log", ACCENT_BLUE,
         "Per-query JSON log written to debug_log_dir:\n• query_analysis fields\n"
         "• retrieval_query string\n• top docs + sections + scores\n"
         "• is_answerable flag\n• full LLM prompt payload\n• final answer"),
        ("Ingestion Progress", ACCENT_PURP,
         "Real-time progress callbacks:\n• processed_files / total_files\n"
         "• current_file name\n• percent complete\n• Phase: extract|summarise|embed|index\n"
         "• Streamed to UI via health poll (1.5s interval during indexing)"),
    ]

    for i, (title, col, body) in enumerate(metric_cards):
        bx = 0.35 + i * 3.24
        add_rounded_rect(s, bx, 1.38, 3.10, 3.80, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.2))
        add_rect(s, bx, 1.38, 3.10, 0.07, fill_color=col)
        add_text_box(s, title, bx + 0.14, 1.48, 2.82, 0.28,
                     font_size=Pt(10.5), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, bx + 0.14, 1.84, 2.82, 3.25,
                     font_size=Pt(9.5), color=LIGHT_GREY, font_name="Segoe UI")

    # Bottom: key engineering decisions
    add_text_box(s, "Key Engineering Decisions & Rationale", 0.40, 5.32, 12, 0.28,
                 font_size=Pt(11), bold=True, color=ACCENT_TEAL,
                 font_name="Segoe UI Semibold")
    decisions = [
        ("Embedding on summary, not raw text",
         "Raw text includes boilerplate, headers, formatting. Summaries are semantically dense — "
         "cosine similarity on summaries aligns with user intent better than noisy raw text."),
        ("Concurrent search with ThreadPoolExecutor",
         "All three search types run in parallel. Total latency = max(k, s, v) not k+s+v. "
         "Individual failures are absorbed — degraded results rather than errors."),
        ("In-memory kNN matrix kept for fallback",
         "When OS is unavailable or hybrid_search_enabled=False, the NumPy kNN matrix is "
         "used. corpus.rebuild_indexes() skips empty-embedding docs (OS-cached) to avoid np.vstack crashes."),
    ]
    for i, (title, body) in enumerate(decisions):
        by = 5.68 + i * 0.52
        add_rounded_rect(s, 0.35, by, 12.63, 0.46, fill_color=CARD_BG,
                         line_color=DIVIDER, line_width=Pt(0.5))
        add_rect(s, 0.35, by, 0.05, 0.46, fill_color=ACCENT_TEAL)
        add_text_box(s, title, 0.52, by + 0.03, 3.30, 0.36,
                     font_size=Pt(9.5), bold=True, color=WHITE,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, 3.88, by + 0.03, 8.90, 0.36,
                     font_size=Pt(9), color=LIGHT_GREY, font_name="Segoe UI")
    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 – Summary / Closing
# ═══════════════════════════════════════════════════════════════════════════════
def slide_summary():
    s = prs.slides.add_slide(blank_layout)
    slide_background(s)
    add_rect(s, 0, 0, 0.06, 7.50, fill_color=ACCENT_TEAL)

    add_slide_label(s, "SUMMARY")
    add_text_box(s, "What We Built", 0.35, 0.65, 12, 0.80,
                 font_size=Pt(36), bold=True, color=WHITE,
                 font_name="Segoe UI Semibold")
    h_divider(s, 1.55, color=ACCENT_TEAL)

    summary_items = [
        (ACCENT_TEAL,  "Hybrid Retrieval",
         "3-channel search (BM25 + MLT + kNN) blended with min-max normalisation and configurable weights per domain"),
        (ACCENT_BLUE,  "Decoupled Ingestion & Retrieval",
         "Background async ingestion; server is 'ready' immediately; retrieval goes straight to OpenSearch"),
        (ACCENT_PURP,  "FAQ Fast-Path",
         "Pre-generated Q/A index with kNN lookup bypasses RAG entirely for high-frequency queries"),
        (ACCENT_AMBER, "Domain Optimisation",
         "Single domain_type switch changes weights, OCR, token budgets, prompts — no code changes"),
        (ACCENT_GREEN, "Permission-Native Security",
         "user_ids array on every OS document; pre-filter at query time; zero post-hoc filtering"),
        (ACCENT_TEAL,  "Full Conversation Memory",
         "recent_messages for LLM context; display_messages accumulated in OS; auto-summarisation at 8 turns"),
        (ACCENT_BLUE,  "Pluggable Backends",
         "VectorStore abstract port; swap opensearch → pinecone / weaviate / pgvector with zero other changes"),
        (ACCENT_PURP,  "Observable & Tunable",
         "Usage tracker, retrieval logs, health endpoint, accuracy profiles, runtime cost profiles"),
    ]

    for i, (col, title, body) in enumerate(summary_items):
        row = i // 2
        c   = i % 2
        bx  = 0.35 + c * 6.45
        by  = 1.68 + row * 1.24
        add_rounded_rect(s, bx, by, 6.25, 1.10, fill_color=CARD_BG,
                         line_color=col, line_width=Pt(1.0))
        add_rect(s, bx, by, 0.06, 1.10, fill_color=col)
        add_text_box(s, title, bx + 0.20, by + 0.08, 5.90, 0.28,
                     font_size=Pt(11), bold=True, color=col,
                     font_name="Segoe UI Semibold")
        add_text_box(s, body, bx + 0.20, by + 0.42, 5.90, 0.60,
                     font_size=Pt(9.5), color=LIGHT_GREY, font_name="Segoe UI")

    footer(s)


# ═══════════════════════════════════════════════════════════════════════════════
# Build all slides
# ═══════════════════════════════════════════════════════════════════════════════
slide_cover()
slide_concept()
slide_architecture()
slide_ingestion_flow()
slide_retrieval_flow()
slide_index_schema()
slide_score_blending()
slide_access_control()
slide_domain_optimisation()
slide_conversation()
slide_faq()
slide_prompt()
slide_deployment()
slide_metrics()
slide_summary()

out = r"d:\policy-mgmt\policygpt-poc\PolicyGPT_HybridSearch_Engineering.pptx"
prs.save(out)
print(f"Saved: {out}")
